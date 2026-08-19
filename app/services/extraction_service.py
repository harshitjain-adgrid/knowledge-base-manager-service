"""
Text extraction for every file format the knowledge base accepts.

Every extractor turns raw bytes into plain text plus a little metadata. Once a
file is text, the rest of the pipeline (chunking → embedding → pgvector) is
identical regardless of where it came from, so adding a format only ever means
adding an extractor here and registering its extension.
"""

import csv
import datetime
import io
import json
import logging
import re
from dataclasses import dataclass, field

import yaml

from app.services.pdf_service import extract_text_from_pdf

logger = logging.getLogger(__name__)


@dataclass
class ExtractedDocument:
    """Normalised result of extracting text from an uploaded file."""

    text: str
    file_name: str
    file_size: int
    source_format: str
    metadata: dict = field(default_factory=dict)


# ── Helpers ──

def _decode(data: bytes, file_name: str) -> str:
    """Decode bytes to text, tolerating the usual Windows encodings."""
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    # latin-1 never fails, so this is unreachable in practice
    raise ValueError(f"Could not decode '{file_name}' as text.")


_FRONTMATTER_RE = re.compile(r"^---\s*\n(.*?)\n---\s*\n", re.DOTALL)


def _split_frontmatter(text: str) -> tuple[dict, str]:
    """
    Pull YAML front matter off the top of a markdown file.

    Returns (frontmatter_dict, body). Malformed front matter is left in the body
    rather than raising — a bad YAML block should not block an upload.
    """
    match = _FRONTMATTER_RE.match(text)
    if not match:
        return {}, text

    try:
        parsed = yaml.safe_load(match.group(1))
    except yaml.YAMLError as e:
        logger.warning(f"Ignoring malformed YAML front matter: {e}")
        return {}, text

    if not isinstance(parsed, dict):
        return {}, text

    return parsed, text[match.end():]


# ── Extractors ──

def _extract_pdf(data: bytes, file_name: str) -> ExtractedDocument:
    extraction = extract_text_from_pdf(data, file_name)
    return ExtractedDocument(
        text=extraction.text,
        file_name=extraction.file_name,
        file_size=extraction.file_size,
        source_format="pdf",
        metadata={"page_count": extraction.page_count},
    )


def _extract_docx(data: bytes, file_name: str) -> ExtractedDocument:
    from docx import Document as DocxDocument

    try:
        document = DocxDocument(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"Failed to read Word file '{file_name}': {e}")

    blocks: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        # Promote Word headings to markdown so the chunker can see structure
        style = (paragraph.style.name or "") if paragraph.style else ""
        if style.startswith("Heading"):
            level = "".join(c for c in style if c.isdigit())
            blocks.append(f"{'#' * min(int(level or 1), 6)} {text}")
        else:
            blocks.append(text)

    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append("\n".join(rows))

    text = "\n\n".join(blocks)
    if not text.strip():
        raise ValueError(f"Word file '{file_name}' contains no extractable text.")

    return ExtractedDocument(
        text=text,
        file_name=file_name,
        file_size=len(data),
        source_format="docx",
        metadata={
            "paragraph_count": len(document.paragraphs),
            "table_count": len(document.tables),
        },
    )


def _extract_txt(data: bytes, file_name: str) -> ExtractedDocument:
    text = _decode(data, file_name)
    if not text.strip():
        raise ValueError(f"File '{file_name}' is empty.")
    return ExtractedDocument(
        text=text,
        file_name=file_name,
        file_size=len(data),
        source_format="txt",
    )


def _extract_markdown(data: bytes, file_name: str) -> ExtractedDocument:
    raw = _decode(data, file_name)
    frontmatter, body = _split_frontmatter(raw)

    if not body.strip():
        raise ValueError(f"Markdown file '{file_name}' has no body content.")

    metadata: dict = {}
    if frontmatter:
        # Keep it namespaced so it can never collide with our own metadata keys
        metadata["frontmatter"] = frontmatter

    return ExtractedDocument(
        text=body.strip(),
        file_name=file_name,
        file_size=len(data),
        source_format="md",
        metadata=metadata,
    )


_HTML_STRIP_TAGS = ("script", "style", "nav", "footer", "header", "noscript", "svg")


def _extract_html(data: bytes, file_name: str) -> ExtractedDocument:
    from bs4 import BeautifulSoup

    raw = _decode(data, file_name)
    soup = BeautifulSoup(raw, "html.parser")

    for tag in soup(_HTML_STRIP_TAGS):
        tag.decompose()

    # Promote headings to markdown before flattening, so structure survives
    for level in range(1, 7):
        for heading in soup.find_all(f"h{level}"):
            heading.replace_with(f"\n\n{'#' * level} {heading.get_text(strip=True)}\n")

    title_tag = soup.find("title")
    page_title = title_tag.get_text(strip=True) if title_tag else None

    text = soup.get_text(separator="\n")
    # Collapse the ocean of blank lines HTML flattening produces
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = "\n".join(line.strip() for line in text.splitlines())
    text = re.sub(r"\n{3,}", "\n\n", text).strip()

    if not text:
        raise ValueError(f"HTML file '{file_name}' contains no extractable text.")

    metadata = {"page_title": page_title} if page_title else {}
    return ExtractedDocument(
        text=text,
        file_name=file_name,
        file_size=len(data),
        source_format="html",
        metadata=metadata,
    )


def _extract_pptx(data: bytes, file_name: str) -> ExtractedDocument:
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"Failed to read PowerPoint file '{file_name}': {e}")

    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                lines.append(text)
        if lines:
            slides.append(f"## Slide {index}\n\n" + "\n".join(lines))

    text = "\n\n".join(slides)
    if not text.strip():
        raise ValueError(
            f"PowerPoint file '{file_name}' contains no extractable text. "
            f"Slides made only of images are not supported."
        )

    return ExtractedDocument(
        text=text,
        file_name=file_name,
        file_size=len(data),
        source_format="pptx",
        metadata={"slide_count": len(presentation.slides)},
    )


def _rows_to_text(header: list[str], rows: list[list[str]], label: str) -> list[str]:
    """
    Render tabular rows one record per block, repeating the column names.

    Repeating headers is deliberate: a chunk that reads 'valid_until: 2026-03-01'
    is retrievable, whereas a bare cell value is not.
    """
    blocks = []
    for row in rows:
        pairs = [
            f"{col}: {val}"
            for col, val in zip(header, row)
            if str(val).strip() not in ("", "None")
        ]
        if pairs:
            blocks.append(f"{label}\n" + "\n".join(pairs))
    return blocks


def _extract_csv(data: bytes, file_name: str) -> ExtractedDocument:
    raw = _decode(data, file_name)
    reader = csv.reader(io.StringIO(raw))

    try:
        all_rows = [row for row in reader if any(cell.strip() for cell in row)]
    except csv.Error as e:
        raise ValueError(f"Failed to parse CSV '{file_name}': {e}")

    if not all_rows:
        raise ValueError(f"CSV file '{file_name}' is empty.")

    header, rows = all_rows[0], all_rows[1:]
    blocks = _rows_to_text(header, rows, f"Row from {file_name}:")

    if not blocks:
        raise ValueError(f"CSV file '{file_name}' has a header but no data rows.")

    return ExtractedDocument(
        text="\n\n".join(blocks),
        file_name=file_name,
        file_size=len(data),
        source_format="csv",
        metadata={"row_count": len(rows), "columns": header},
    )


def _extract_xlsx(data: bytes, file_name: str) -> ExtractedDocument:
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as e:
        raise ValueError(f"Failed to read Excel file '{file_name}': {e}")

    blocks: list[str] = []
    sheet_names: list[str] = []

    for sheet in workbook.worksheets:
        values = [
            ["" if cell is None else str(cell) for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        values = [row for row in values if any(cell.strip() for cell in row)]
        if not values:
            continue

        sheet_names.append(sheet.title)
        header, rows = values[0], values[1:]
        blocks.extend(_rows_to_text(header, rows, f"Row from sheet '{sheet.title}':"))

    workbook.close()

    if not blocks:
        raise ValueError(f"Excel file '{file_name}' contains no data rows.")

    return ExtractedDocument(
        text="\n\n".join(blocks),
        file_name=file_name,
        file_size=len(data),
        source_format="xlsx",
        metadata={"sheets": sheet_names},
    )


def _extract_json(data: bytes, file_name: str) -> ExtractedDocument:
    raw = _decode(data, file_name)
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError as e:
        raise ValueError(f"Failed to parse JSON '{file_name}': {e}")

    return ExtractedDocument(
        # Re-serialise so the stored content is consistently formatted, and so
        # the api_definition chunker gets valid JSON regardless of input styling
        text=json.dumps(parsed, indent=2, ensure_ascii=False),
        file_name=file_name,
        file_size=len(data),
        source_format="json",
    )


# ── Registry ──

EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".md": _extract_markdown,
    ".markdown": _extract_markdown,
    ".txt": _extract_txt,
    ".text": _extract_txt,
    ".log": _extract_txt,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".pptx": _extract_pptx,
    ".csv": _extract_csv,
    ".xlsx": _extract_xlsx,
    ".json": _extract_json,
}

SUPPORTED_EXTENSIONS = sorted(EXTRACTORS.keys())

# Formats where a single record per chunk beats prose chunking, and where the
# admin should be warned that a spreadsheet is usually a poor knowledge source.
TABULAR_FORMATS = {"csv", "xlsx"}


def get_extension(file_name: str) -> str:
    """Lowercased extension including the leading dot, or '' if there is none."""
    _, dot, ext = file_name.rpartition(".")
    return f".{ext.lower()}" if dot else ""


def is_supported(file_name: str) -> bool:
    return get_extension(file_name) in EXTRACTORS


def extract(file_content: bytes, file_name: str) -> ExtractedDocument:
    """
    Extract text from an uploaded file, dispatching on its extension.

    Raises:
        ValueError: unsupported extension, unreadable file, or no text found.
    """
    if not file_name:
        raise ValueError("Uploaded file has no name, so its format is unknown.")

    extension = get_extension(file_name)
    extractor = EXTRACTORS.get(extension)
    if extractor is None:
        raise ValueError(
            f"Unsupported file type '{extension or file_name}'. "
            f"Supported types: {', '.join(SUPPORTED_EXTENSIONS)}."
        )

    if not file_content:
        raise ValueError(f"Uploaded file '{file_name}' is empty.")

    result = extractor(file_content, file_name)
    logger.info(
        f"Extracted {len(result.text)} characters from '{file_name}' "
        f"({result.source_format})."
    )
    return result


# Front matter keys that map onto real document fields rather than free-form
# metadata. Everything else in the front matter is kept as metadata verbatim.
_FM_TITLE_KEYS = ("title",)
_FM_TYPE_KEYS = ("type", "doc_type")
_DOC_TYPE_RE = re.compile(r"^[a-z0-9_-]+$")


def _jsonable(value: object) -> object:
    """
    Coerce a YAML value into something JSONB can store.

    PyYAML resolves `last_reviewed: 2026-08-18` to a datetime.date, and dates,
    times and sets all fail json.dumps — which would 500 the upload rather than
    reject the field. Dates become ISO strings; anything else exotic becomes its
    string form. Containers are converted recursively.
    """
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, (datetime.datetime, datetime.date, datetime.time)):
        return value.isoformat()
    if isinstance(value, dict):
        return {str(k): _jsonable(v) for k, v in value.items()}
    if isinstance(value, (list, tuple, set, frozenset)):
        return [_jsonable(v) for v in value]
    return str(value)


def _clean_doc_type(value: object) -> str | None:
    """Coerce a front matter type into the doc_type vocabulary, or reject it."""
    if not isinstance(value, str):
        return None
    candidate = value.strip().lower().replace(" ", "_")
    if not candidate or not _DOC_TYPE_RE.match(candidate) or len(candidate) > 64:
        logger.warning(f"Ignoring unusable front matter type {value!r}.")
        return None
    return candidate


def promote_frontmatter(extraction: ExtractedDocument) -> dict:
    """
    Split a document's front matter into real fields and leftover metadata.

    `title` and `type` become document columns; everything else stays metadata
    so it can be filtered on later. This is what makes the authoring format
    meaningful — a team writing `type: capability` gets a document that is
    actually typed, not a comment nobody reads.

    Returns {"title": str | None, "doc_type": str | None, "metadata": dict}.
    """
    frontmatter = extraction.metadata.get("frontmatter")
    if not isinstance(frontmatter, dict):
        return {"title": None, "doc_type": None, "metadata": {}}

    title = None
    doc_type = None
    leftover: dict = {}

    for key, value in frontmatter.items():
        lowered = str(key).strip().lower()
        if lowered in _FM_TITLE_KEYS and isinstance(value, str) and value.strip():
            title = value.strip()
        elif lowered in _FM_TYPE_KEYS:
            doc_type = _clean_doc_type(value)
            if doc_type is None and value is not None:
                leftover[lowered] = _jsonable(value)
        else:
            leftover[lowered] = _jsonable(value)

    return {"title": title, "doc_type": doc_type, "metadata": leftover}


def suggested_doc_type(source_format: str) -> str:
    """
    Pick the chunking strategy for a freshly extracted file.

    JSON is assumed to be an API definition; everything else is prose.
    """
    return "api_definition" if source_format == "json" else "text"
